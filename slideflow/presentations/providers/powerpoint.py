"""PowerPoint presentation provider for local PPTX output.

The provider renders Slideflow presentations into native ``.pptx`` files using
``python-pptx``. It mirrors the presentation provider interface used by cloud
providers while keeping generated chart images in memory until they are inserted
into the destination deck.
"""

from __future__ import annotations

import hashlib
import os
import re
import stat
import uuid
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, Iterable, List, Literal, Optional, Tuple
from urllib.request import urlopen

from pydantic import Field, field_validator

from slideflow.constants import GoogleSlides
from slideflow.presentations.providers.base import (
    PresentationProvider,
    PresentationProviderConfig,
)
from slideflow.utilities.exceptions import ConfigurationError, RenderingError
from slideflow.utilities.logging import get_logger

try:  # pragma: no cover - exercised in environments without the optional extra.
    from pptx import Presentation as PptxPresentation
except ImportError:  # pragma: no cover
    PptxPresentation = None  # type: ignore[assignment]

logger = get_logger(__name__)

EMU_PER_POINT = 12700
MEMORY_IMAGE_URL_PREFIX = "slideflow-pptx://chart/"
AUXILIARY_STEM_MAX_LENGTH = 80


def _require_python_pptx() -> Any:
    if PptxPresentation is None:
        raise ConfigurationError(
            "PowerPoint provider requires the optional dependency "
            "'python-pptx'. Install with slideflow-presentations[powerpoint]."
        )
    return PptxPresentation


def _normalize_path(value: Path | str) -> Path:
    return Path(value).expanduser()


def _nearest_existing_parent(path: Path) -> Optional[Path]:
    current = path if path.exists() else path.parent
    while not current.exists():
        if current == current.parent:
            return None
        current = current.parent
    return current


def _safe_output_stem(name: str) -> str:
    stem = re.sub(r"[^A-Za-z0-9._ -]+", "_", name).strip(" ._")
    stem = re.sub(r"\s+", " ", stem)
    return stem or "presentation"


def _points_to_emu(value: float) -> int:
    return int(round(value * EMU_PER_POINT))


def _emu_to_points(value: int) -> int:
    return int(round(value / EMU_PER_POINT))


class PowerPointProviderConfig(PresentationProviderConfig):
    """Configuration for the local PowerPoint PPTX provider."""

    provider_type: Literal["powerpoint"] = "powerpoint"
    template_path: Path = Field(..., description="Path to the source .pptx template")
    output_dir: Path = Field(
        Path("."),
        description="Directory where generated .pptx files are written.",
    )
    slide_id_mode: Literal["auto", "index", "native"] = Field(
        "auto",
        description=(
            "How presentation.slides[].id resolves slides: one-based index, "
            "native PowerPoint slide_id, or auto (index then native)."
        ),
    )
    read_only_template: bool = Field(
        True,
        description="Prevent writes to the configured template_path.",
    )
    file_collision_strategy: Literal["fail", "overwrite", "suffix"] = Field(
        "fail",
        description="Behavior when the generated .pptx path already exists.",
    )
    strict_cleanup: bool = Field(
        False,
        description="Fail rendering when temporary in-memory chart cleanup fails.",
    )
    share_with: List[str] = Field(
        default_factory=list,
        description="Accepted for provider compatibility; sharing is a no-op.",
    )
    share_role: str = Field(
        GoogleSlides.PERMISSION_READER,
        description="Accepted for provider compatibility; sharing is a no-op.",
    )

    @field_validator("template_path")
    @classmethod
    def _validate_template_path(cls, value: Path | str) -> Path:
        path = _normalize_path(value)
        if path.suffix.lower() != ".pptx":
            raise ValueError("template_path must point to a .pptx file")
        return path

    @field_validator("output_dir")
    @classmethod
    def _validate_output_dir(cls, value: Path | str) -> Path:
        return _normalize_path(value)


class PowerPointProvider(PresentationProvider):
    """Local PowerPoint provider that renders native ``.pptx`` artifacts."""

    def __init__(self, config: PowerPointProviderConfig):
        super().__init__(config)
        self.config: PowerPointProviderConfig = config
        self._pptx_factory = _require_python_pptx()
        self._presentations: Dict[str, Any] = {}
        self._presentation_paths: Dict[str, Path] = {}
        self._presentation_locks: Dict[str, Path] = {}
        self._chart_images: Dict[str, bytes] = {}

    def run_preflight_checks(self) -> List[Tuple[str, bool, str]]:
        """Run local PPTX provider preflight checks for doctor/build."""
        checks: List[Tuple[str, bool, str]] = [
            (
                "python_pptx_import",
                PptxPresentation is not None,
                (
                    "python-pptx import succeeded"
                    if PptxPresentation is not None
                    else "Install slideflow-presentations[powerpoint]"
                ),
            )
        ]

        template_path = self.config.template_path
        checks.append(
            (
                "template_path_exists",
                template_path.is_file(),
                (
                    f"Found template at {template_path}"
                    if template_path.is_file()
                    else f"Template file not found: {template_path}"
                ),
            )
        )
        checks.append(
            (
                "template_path_is_pptx",
                template_path.suffix.lower() == ".pptx",
                f"Template suffix is '{template_path.suffix or '<none>'}'",
            )
        )

        output_dir = self.config.output_dir
        output_parent = _nearest_existing_parent(output_dir)
        output_parent_ok = (
            output_parent is not None
            and output_parent.is_dir()
            and os.access(output_parent, os.W_OK)
        )
        checks.append(
            (
                "output_dir_writable",
                output_parent_ok,
                (
                    f"Output directory is available: {output_dir}"
                    if output_dir.exists()
                    else f"Output directory can be created under: {output_parent}"
                ),
            )
        )
        return checks

    def create_presentation(self, name: str, template_id: Optional[str] = None) -> str:
        """Create a local PPTX artifact from the configured template."""
        template_path = _normalize_path(template_id or self.config.template_path)
        if template_path.suffix.lower() != ".pptx":
            raise ConfigurationError("PowerPoint template must be a .pptx file")
        if not template_path.is_file():
            raise ConfigurationError(f"PowerPoint template not found: {template_path}")

        output_path, output_lock_path = self._resolve_output_path(name)
        if (
            output_path.resolve() == template_path.resolve()
            and self.config.read_only_template
        ):
            if output_lock_path is not None:
                self._release_output_lock(output_lock_path)
            raise ConfigurationError(
                "PowerPoint output path resolves to template_path while "
                "read_only_template is true"
            )

        try:
            presentation = self._pptx_factory(str(template_path))
        except BaseException:
            if output_lock_path is not None:
                self._release_output_lock(output_lock_path)
            raise

        presentation_id = str(output_path)
        self._presentations[presentation_id] = presentation
        self._presentation_paths[presentation_id] = output_path
        if output_lock_path is not None:
            self._presentation_locks[presentation_id] = output_lock_path
        return presentation_id

    def upload_chart_image(
        self, presentation_id: str, image_data: bytes, filename: str
    ) -> Tuple[str, str]:
        """Store chart image bytes in memory for later insertion."""
        self._get_presentation(presentation_id)
        token = uuid.uuid4().hex
        image_url = f"{MEMORY_IMAGE_URL_PREFIX}{token}/{Path(filename).name}"
        self._chart_images[token] = bytes(image_data)
        return image_url, token

    def insert_chart_to_slide(
        self,
        presentation_id: str,
        slide_id: str,
        image_url: str,
        x: float,
        y: float,
        width: float,
        height: float,
    ) -> None:
        """Insert a chart image into a PowerPoint slide."""
        if width <= 0 or height <= 0:
            raise RenderingError("Chart width and height must be greater than zero")

        slide = self._resolve_slide(presentation_id, slide_id)
        image_stream = self._open_image_stream(image_url)
        slide.shapes.add_picture(
            image_stream,
            _points_to_emu(x),
            _points_to_emu(y),
            width=_points_to_emu(width),
            height=_points_to_emu(height),
        )

    def replace_text_in_slide(
        self, presentation_id: str, slide_id: str, placeholder: str, replacement: str
    ) -> int:
        """Replace placeholder text in text frames and table cells on one slide."""
        if not placeholder:
            return 0

        slide = self._resolve_slide(presentation_id, slide_id)
        replacements = 0
        for text_frame in self._iter_text_frames(slide.shapes):
            replacements += self._replace_text_frame(
                text_frame=text_frame,
                placeholder=placeholder,
                replacement=replacement,
            )
        return replacements

    def share_presentation(
        self,
        presentation_id: str,
        emails: List[str],
        role: str = GoogleSlides.PERMISSION_READER,
    ) -> None:
        """No-op sharing hook for local PPTX output."""
        self._get_presentation(presentation_id)
        if emails:
            logger.warning(
                "PowerPoint provider does not support sharing; generated PPTX "
                "remains local (requested role=%s, recipients=%s)",
                role,
                ", ".join(emails),
            )

    def get_presentation_url(self, presentation_id: str) -> str:
        """Return a file URL for the generated PPTX artifact."""
        output_path = self._presentation_paths.get(
            presentation_id, Path(presentation_id)
        )
        return output_path.resolve().as_uri()

    def get_presentation_page_size(
        self, presentation_id: str
    ) -> Optional[Tuple[int, int]]:
        """Return slide dimensions in points."""
        presentation = self._get_presentation(presentation_id)
        return _emu_to_points(int(presentation.slide_width)), _emu_to_points(
            int(presentation.slide_height)
        )

    def finalize_presentation(self, presentation_id: str) -> None:
        """Save the generated PPTX artifact to disk."""
        presentation = self._get_presentation(presentation_id)
        output_path = self._presentation_paths[presentation_id]
        output_lock_path = self._presentation_locks.get(presentation_id)
        temp_path = self._temp_path_for(output_path)
        existing_mode = self._existing_file_mode(output_path)

        try:
            if output_lock_path is not None and output_path.exists():
                raise ConfigurationError(
                    f"PowerPoint output appeared after reservation: {output_path}"
                )
            if existing_mode is not None:
                with temp_path.open("xb"):
                    pass
                temp_path.chmod(existing_mode)
            presentation.save(str(temp_path))
            if existing_mode is not None:
                temp_path.chmod(existing_mode)
            temp_path.replace(output_path)
        finally:
            temp_path.unlink(missing_ok=True)
            self._release_presentation_lock(presentation_id)

    def abort_presentation(self, presentation_id: str) -> None:
        """Release a reserved output path after an unsuccessful render."""
        self._release_presentation_lock(presentation_id)
        self._presentations.pop(presentation_id, None)
        self._presentation_paths.pop(presentation_id, None)

    def render_citations(
        self,
        presentation_id: str,
        citations_by_scope: Dict[str, List[Dict[str, Any]]],
        location: str,
    ) -> None:
        """Render source citations into PowerPoint speaker notes when possible.

        ``python-pptx`` does not expose a stable public API for notes pages, so
        PPTX citations are intentionally not rendered in V1.
        """
        self._get_presentation(presentation_id)
        if citations_by_scope:
            logger.warning(
                "PowerPoint provider does not render citations yet "
                "(location=%s, scopes=%d)",
                location,
                len(citations_by_scope),
            )

    def delete_chart_image(self, file_id: str) -> None:
        """Delete an in-memory chart image token."""
        self._chart_images.pop(file_id, None)

    def _resolve_output_path(self, name: str) -> Tuple[Path, Optional[Path]]:
        output_dir = self.config.output_dir
        output_path = output_dir / f"{_safe_output_stem(name)}.pptx"
        output_path.parent.mkdir(parents=True, exist_ok=True)

        if self.config.file_collision_strategy == "overwrite":
            if (
                output_path.resolve() == self.config.template_path.resolve()
                and self.config.read_only_template
            ):
                raise ConfigurationError(
                    "Refusing to overwrite template_path while "
                    "read_only_template is true"
                )
            return output_path, None

        try:
            return output_path, self._reserve_output_lock(output_path)
        except FileExistsError:
            if self.config.file_collision_strategy == "fail":
                raise ConfigurationError(
                    f"PowerPoint output already exists: {output_path}"
                ) from None

        return self._suffix_output_path(output_path)

    @staticmethod
    def _suffix_output_path(output_path: Path) -> Tuple[Path, Path]:
        stem = output_path.stem
        suffix = output_path.suffix
        parent = output_path.parent
        for index in range(1, 10_000):
            candidate = parent / f"{stem}-{index}{suffix}"
            try:
                return candidate, PowerPointProvider._reserve_output_lock(candidate)
            except FileExistsError:
                continue
        raise ConfigurationError(
            f"Could not find available suffixed output path for {output_path}"
        )

    @staticmethod
    def _lock_path_for(output_path: Path) -> Path:
        return output_path.with_name(
            f"{PowerPointProvider._auxiliary_name_prefix(output_path)}.lock"
        )

    @staticmethod
    def _temp_path_for(output_path: Path) -> Path:
        return output_path.with_name(
            f"{PowerPointProvider._auxiliary_name_prefix(output_path)}."
            f"{uuid.uuid4().hex}.tmp"
        )

    @staticmethod
    def _auxiliary_name_prefix(output_path: Path) -> str:
        stem = output_path.stem[:AUXILIARY_STEM_MAX_LENGTH].rstrip(" .")
        stem = stem or "presentation"
        digest = hashlib.sha256(output_path.name.encode("utf-8")).hexdigest()[:16]
        return f".{stem}.{digest}"

    @staticmethod
    def _existing_file_mode(output_path: Path) -> Optional[int]:
        try:
            return stat.S_IMODE(output_path.stat().st_mode)
        except FileNotFoundError:
            return None

    @staticmethod
    def _reserve_output_lock(output_path: Path) -> Path:
        if output_path.exists():
            raise FileExistsError(output_path)

        lock_path = PowerPointProvider._lock_path_for(output_path)
        try:
            with lock_path.open("xb"):
                pass
        except FileExistsError:
            raise
        except OSError as error:
            raise ConfigurationError(
                f"Could not reserve PowerPoint output path {output_path}: {error}"
            ) from error

        if output_path.exists():
            PowerPointProvider._release_output_lock(lock_path)
            raise FileExistsError(output_path)
        return lock_path

    def _release_presentation_lock(self, presentation_id: str) -> None:
        lock_path = self._presentation_locks.pop(presentation_id, None)
        if lock_path is not None:
            self._release_output_lock(lock_path)

    @staticmethod
    def _release_output_lock(lock_path: Path) -> None:
        try:
            lock_path.unlink(missing_ok=True)
        except OSError as error:
            logger.warning(
                "Could not release PowerPoint output lock %s: %s", lock_path, error
            )

    def _get_presentation(self, presentation_id: str) -> Any:
        try:
            return self._presentations[presentation_id]
        except KeyError as error:
            raise RenderingError(
                f"Unknown PowerPoint presentation id: {presentation_id}"
            ) from error

    def _resolve_slide(self, presentation_id: str, slide_id: str) -> Any:
        presentation = self._get_presentation(presentation_id)
        slides = list(presentation.slides)

        if self.config.slide_id_mode in {"auto", "index"}:
            slide = self._slide_by_one_based_index(slides, slide_id)
            if slide is not None:
                return slide
            if self.config.slide_id_mode == "index":
                raise RenderingError(
                    f"PowerPoint slide index '{slide_id}' is not present "
                    f"(slide count: {len(slides)})"
                )

        if self.config.slide_id_mode in {"auto", "native"}:
            for slide in slides:
                if str(getattr(slide, "slide_id", "")) == str(slide_id):
                    return slide

        raise RenderingError(
            f"PowerPoint slide id '{slide_id}' was not found using "
            f"slide_id_mode='{self.config.slide_id_mode}'"
        )

    @staticmethod
    def _slide_by_one_based_index(slides: List[Any], slide_id: str) -> Optional[Any]:
        try:
            index = int(str(slide_id).strip())
        except ValueError:
            return None
        if index < 1 or index > len(slides):
            return None
        return slides[index - 1]

    def _open_image_stream(self, image_url: str) -> BytesIO:
        token = self._extract_memory_image_token(image_url)
        if token is not None:
            try:
                return BytesIO(self._chart_images[token])
            except KeyError as error:
                raise RenderingError(
                    f"PowerPoint chart image token is no longer available: {token}"
                ) from error

        if image_url.startswith(("http://", "https://")):
            try:
                with urlopen(image_url, timeout=15) as response:  # nosec B310
                    return BytesIO(response.read())
            except Exception as error:
                raise RenderingError(
                    f"Failed to download chart image for PowerPoint insertion: {error}"
                ) from error

        local_path = Path(image_url).expanduser()
        if local_path.is_file():
            return BytesIO(local_path.read_bytes())

        raise RenderingError(
            "PowerPoint chart image URL must be an uploaded memory token, "
            "HTTP(S) URL, or readable local file path"
        )

    @staticmethod
    def _extract_memory_image_token(image_url: str) -> Optional[str]:
        if not image_url.startswith(MEMORY_IMAGE_URL_PREFIX):
            return None
        remainder = image_url[len(MEMORY_IMAGE_URL_PREFIX) :]
        return remainder.split("/", 1)[0] or None

    def _iter_text_frames(self, shapes: Iterable[Any]) -> Iterable[Any]:
        for shape in shapes:
            if getattr(shape, "has_text_frame", False):
                yield shape.text_frame

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame

            child_shapes = getattr(shape, "shapes", None)
            if child_shapes is not None:
                yield from self._iter_text_frames(child_shapes)

    @staticmethod
    def _replace_text_frame(text_frame: Any, placeholder: str, replacement: str) -> int:
        replacements = 0
        for paragraph in text_frame.paragraphs:
            original_paragraph_text = paragraph.text
            original_count = original_paragraph_text.count(placeholder)
            paragraph_replacements = 0
            for run in paragraph.runs:
                if placeholder not in run.text:
                    continue
                paragraph_replacements += run.text.count(placeholder)
                run.text = run.text.replace(placeholder, replacement)

            replacements += paragraph_replacements
            if original_count > paragraph_replacements:
                replacements += original_count - paragraph_replacements
                paragraph.text = original_paragraph_text.replace(
                    placeholder, replacement
                )

        return replacements
